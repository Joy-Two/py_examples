from pptx import Presentation
from pptx.util import Inches,Pt


'''
    learning from https://python-pptx.readthedocs.io/en/latest/user/quickstart.html#add-table-example
'''

prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

# shapes.title.text = ''

rows = 2
cols = 5
# left = top = Inches(2.0)
top = Inches(0)
left = Inches(0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

for i in range(rows):
    for j in range(cols):
        table.cell(i,j).text_frame.paragraphs[0].font.size = Pt(11)

# set column widths
table.columns[0].width = Inches(1.0)
table.columns[1].width = Inches(1.0)
table.columns[2].width = Inches(1.0)
table.columns[3].width = Inches(1.0)
table.columns[4].width = Inches(1.0)

# write column headings
table.cell(0, 0).text = 'Group'
table.cell(0, 1).text = 'Total Projects'
table.cell(0, 2).text = 'Items1'
table.cell(0, 3).text = 'Items2'
table.cell(0, 4).text = 'Items3'

# write body cells
table.cell(1, 0).text = 'Content Eng/Pub'
table.cell(1, 1).text = '5'
table.cell(1, 2).text = '100%'
table.cell(1, 2).fill.background ='transparent'
table.cell(1, 3).text = '100%'
table.cell(1, 4).text = '100%'

import os
file = 'test_table.pptx'
os.remove(file)
prs.save(file)